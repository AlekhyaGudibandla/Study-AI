import os
import fitz # PyMuPDF
import docx
from pptx import Presentation
from celery import Celery
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
from sqlalchemy.orm import Session
from cryptography.hazmat.primitives import hashes

from database import engine
from models import Document

# Use Redis URL from environment or fallback
CELERY_BROKER_URL = os.getenv("CELERY_BROKER_URL", "redis://localhost:6379/0")
CELERY_RESULT_BACKEND = CELERY_BROKER_URL

celery_app = Celery("rag_worker", broker=CELERY_BROKER_URL, backend=CELERY_RESULT_BACKEND)

# Initialize ChromaDB Local Client (saving in the backend directory)
_chroma_client = None
def get_chroma_client():
    global _chroma_client
    if _chroma_client is None:
        print("DEBUG: Initializing ChromaDB client...")
        _chroma_client = chromadb.PersistentClient(path="./chroma_db", settings=Settings(allow_reset=True))
    return _chroma_client

# Embedding model - lazy loaded
_embedder = None
def get_embedder():
    global _embedder
    if _embedder is None:
        print("DEBUG: Loading SentenceTransformer model...")
        _embedder = SentenceTransformer("all-MiniLM-L6-v2")
    return _embedder

def extract_text_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text("text") + "\n"
    return text

def extract_text_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def chunk_text(text, chunk_size=1000, overlap=100):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunks.append(" ".join(words[i:i + chunk_size]))
    return chunks

@celery_app.task(bind=True, max_retries=3)
def process_document_task(self, document_id: str, file_path: str, user_id: str):
    return process_document_logic(document_id, file_path, user_id, task=self)

def process_document_logic(document_id: str, file_path: str, user_id: str, task=None):
    from sqlalchemy.orm import sessionmaker
    SessionLocal = sessionmaker(bind=engine)
    db = SessionLocal()
    
    doc = db.query(Document).filter(Document.id == document_id).first()
    if not doc:
        db.close()
        return

    doc.status = "PROCESSING"
    db.commit()

    try:
        if doc.filename.lower().endswith(".pdf"):
            text = extract_text_pdf(file_path)
        elif doc.filename.lower().endswith((".docx", ".doc")):
            text = extract_text_docx(file_path)
        elif doc.filename.lower().endswith((".pptx", ".ppt")):
            text = extract_text_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {doc.filename}")
            
        chunks = chunk_text(text)
        if not chunks:
            raise ValueError("No text extracted from document")

        embeddings = get_embedder().encode(chunks).tolist()

        # Get or create chroma collection for this specific user
        collection = get_chroma_client().get_or_create_collection(name=f"user_{user_id}")
        
        ids = [f"{document_id}_{i}" for i in range(len(chunks))]
        metadatas = [{"document_id": document_id, "chunk_index": i} for i in range(len(chunks))]
        
        collection.add(
            embeddings=embeddings,
            documents=chunks,
            metadatas=metadatas,
            ids=ids
        )

        doc.status = "SUCCESS"
        db.commit()
    except Exception as e:
        print(f"Error processing document: {str(e)}")
        doc.status = "FAILED"
        doc.error_message = str(e)
        db.commit()
        if task:
            raise task.retry(exc=e, countdown=60)
        else:
            print(f"Background task failed: {e}")
    finally:
        db.close()
        if os.path.exists(file_path):
            os.remove(file_path)
